from PyPDF2 import PdfReader
from docx import Document
import pandas as pd
from pptx import Presentation

def extract_text(file_path):
    """Extract text from various file types, including PDFs, DOCX, CSV, XLSX, MD, PPT, and PPTX."""
    try:
        if file_path.suffix.lower() == ".pdf":
            reader = PdfReader(str(file_path))
            text = "\n".join([page.extract_text() for page in reader.pages])
            return text.strip()
        elif file_path.suffix.lower() == ".docx":
            doc = Document(file_path)
            text = "\n".join([para.text for para in doc.paragraphs])
            return text.strip()
        elif file_path.suffix.lower() == ".csv":
            df = pd.read_csv(file_path)
            return df.to_string(index=False)
        elif file_path.suffix.lower() == ".xlsx":
            df = pd.read_excel(file_path)
            return df.to_string(index=False)
        elif file_path.suffix.lower() in [".ppt", ".pptx"]:
            prs = Presentation(file_path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n".join(text).strip()
        else:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read().strip()
    except Exception as e:
        print(f"Error extracting text from {file_path}: {e}")
        return ""